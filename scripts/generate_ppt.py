"""Generate the EduPulse viva presentation.

Produces EduPulse_Presentation.pptx at the project root.
Widescreen 16:9, blue theme, 24 slides, speaker notes on every content slide.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "assets" / "figures"
SHOT = ROOT / "assets" / "screenshots"
OUT = ROOT / "EduPulse_Presentation.pptx"

# ---- Palette (EduPulse blue) --------------------------------------
BRAND = RGBColor(0x4F, 0x8B, 0xF9)          # #4F8BF9
BRAND_DARK = RGBColor(0x1D, 0x4E, 0xD8)     # #1D4ED8
INK = RGBColor(0x0F, 0x17, 0x2A)            # near-black
MUTED = RGBColor(0x6B, 0x72, 0x80)          # secondary text
BG_LIGHT = RGBColor(0xF3, 0xF6, 0xFC)       # soft blue-white
ACCENT = RGBColor(0xFB, 0xBF, 0x24)         # amber accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------- helpers

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg_bar(slide, colour=BRAND, height=Inches(0.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text,
                *, size=18, bold=False, colour=INK,
                align=PP_ALIGN.LEFT, italic=False,
                font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = colour
    return box


def add_bullets(slide, left, top, width, height, bullets,
                *, size=18, colour=INK, spacing=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {text}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = colour
            run.font.name = "Calibri"
    return box


def add_page_number_footer(slide, index, total):
    add_textbox(slide, Inches(12.0), Inches(7.05),
                Inches(1.2), Inches(0.3),
                f"{index} / {total}",
                size=10, colour=MUTED,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.4), Inches(7.05),
                Inches(4), Inches(0.3),
                "EduPulse — MCA Project Viva",
                size=10, colour=MUTED, italic=True)


def add_speaker_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def fit_image(slide, path, top, max_w, max_h, *, centre=True, left=None):
    """Insert an image scaled to fit inside a (max_w × max_h) box.

    Preserves aspect ratio using PIL to read the image's native size.
    Centres horizontally by default; pass `left=` for manual placement.
    Returns the resulting picture shape.
    """
    from PIL import Image as PILImage
    with PILImage.open(str(path)) as img:
        native_w, native_h = img.size
    ratio = native_w / native_h
    # try width-first fit, then clamp on height
    w_emu = max_w
    h_emu = int(w_emu / ratio)
    if h_emu > max_h:
        h_emu = max_h
        w_emu = int(h_emu * ratio)
    if left is None and centre:
        left = int((SLIDE_W - w_emu) / 2)
    return slide.shapes.add_picture(str(path), left, top,
                                     width=w_emu, height=h_emu)


def title_bar(slide, section_number, title, subtitle=None):
    """Coloured band at the top of every content slide."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.2))
    band.fill.solid(); band.fill.fore_color.rgb = BRAND
    band.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.12),
                Inches(1.5), Inches(0.35),
                f"SLIDE {section_number:02d}",
                size=11, colour=WHITE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.42),
                Inches(12), Inches(0.7),
                title,
                size=28, colour=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.92),
                    Inches(12), Inches(0.3),
                    subtitle,
                    size=13, colour=WHITE, italic=True)


# --------------------------------------------------------------- slides

def build() -> None:
    prs = new_prs()

    # -------- slide 1: Title
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BRAND_DARK
    bg.line.fill.background()
    accent_bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.9), SLIDE_W, Inches(0.08))
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    add_textbox(s, Inches(0.7), Inches(0.5),
                Inches(6), Inches(0.4),
                "⚡ EduPulse", size=22, bold=True,
                colour=WHITE)
    add_textbox(s, Inches(0.7), Inches(3.2),
                Inches(12), Inches(1.2),
                "EduPulse", size=64, bold=True, colour=WHITE)
    add_textbox(s, Inches(0.7), Inches(4.4),
                Inches(12), Inches(1.0),
                "A Business Intelligence Dashboard for Online\n"
                "Learning Engagement Analytics",
                size=26, colour=WHITE)
    add_textbox(s, Inches(0.7), Inches(6.0),
                Inches(12), Inches(0.35),
                "Deepak Kumar Joshi   ·   Enrollment  O24MCA112387",
                size=15, colour=WHITE)
    add_textbox(s, Inches(0.7), Inches(6.35),
                Inches(12), Inches(0.35),
                "Master of Computer Applications  (Online)   ·   "
                "Chandigarh University  (CU Online)",
                size=15, colour=WHITE)
    add_textbox(s, Inches(0.7), Inches(6.70),
                Inches(12), Inches(0.35),
                "Guide:  Vikas Kumar,  Senior Mentor   ·   "
                "Academic Year  2025–2026",
                size=15, colour=WHITE)
    add_speaker_notes(s,
        "Open with a one-sentence pitch: EduPulse is a Business "
        "Intelligence dashboard for online learning engagement "
        "analytics, built on the OULAD dataset. Say your name, "
        "enrollment number, program, guide name, and academic year. "
        "Aim: 30 seconds on this slide.")

    total = 24  # target slide count

    # -------- slide 2: Agenda
    s = blank_slide(prs)
    title_bar(s, 2, "Agenda")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(6),
                [
                    "Background and motivation",
                    "Problem statement — three gaps in current analytics",
                    "Objectives and scope",
                    "Existing systems and where they fall short",
                    "Proposed system — EduPulse architecture",
                    "Technology stack and dataset (OULAD)",
                    "System design — DFDs, star schema, class diagram",
                    "Machine-learning model for dropout prediction",
                    "Dashboard walk-through (six real screenshots)",
                    "Testing, results, and deployment",
                    "Future scope and Q&A",
                ], size=19, spacing=6)
    add_page_number_footer(s, 2, total)
    add_speaker_notes(s,
        "Read the agenda quickly. Signals structure to the panel — "
        "they will trust that you know where you are heading.")

    # -------- slide 3: Background & Motivation
    s = blank_slide(prs)
    title_bar(s, 3, "Background & Motivation",
              "Why online-learning analytics matters")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Online learning is now a primary delivery channel "
                    "— 300M+ registered learners worldwide.",
                    "Every click, video view, forum post, and quiz "
                    "produces structured behavioural data.",
                    "Universities collect the data but rarely act on it "
                    "— dashboards focus on enrolment, not engagement.",
                    "Students at risk of dropping out are typically "
                    "identified too late for intervention to be effective.",
                    "The Open University released OULAD (2014, CC-BY-4.0) "
                    "as the de facto benchmark: 32,593 students, "
                    "10M+ click events.",
                    "EduPulse uses OULAD to demonstrate what a modern "
                    "BI stack + a supervised ML model can do on "
                    "real learning-analytics data.",
                ], size=18)
    add_page_number_footer(s, 3, total)
    add_speaker_notes(s,
        "Emphasise real-world relevance. Mention 300 million learners "
        "and OULAD's 32,593 students. Panel will appreciate that this "
        "is a real dataset, not a toy example.")

    # -------- slide 4: Problem Statement
    s = blank_slide(prs)
    title_bar(s, 4, "Problem Statement",
              "Three operational gaps that motivated EduPulse")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "1. Engagement-visibility gap — administrators "
                    "have no consolidated view of platform-wide "
                    "engagement over time.",
                    "2. Pattern-recognition gap — peaks and troughs, "
                    "content-type popularity, cross-cohort comparisons "
                    "require ad-hoc SQL.",
                    "3. Predictive-intervention gap — no early-warning "
                    "layer to flag students likely to withdraw.",
                    "EduPulse addresses all three in a single deployable "
                    "artefact — free, reproducible, open-source.",
                ], size=20, spacing=10)
    add_page_number_footer(s, 4, total)
    add_speaker_notes(s,
        "Three gaps to remember: visibility, pattern-recognition, "
        "predictive intervention. Say them slowly. If asked, each "
        "corresponds to a dashboard page (Overview, Engagement "
        "Patterns, Dropout Prediction).")

    # -------- slide 5: Objectives
    s = blank_slide(prs)
    title_bar(s, 5, "Objectives")
    add_bullets(s, Inches(0.7), Inches(1.4),
                Inches(12), Inches(6),
                [
                    "Collect real-world learning-analytics dataset "
                    "(OULAD) and document its structure and licence.",
                    "Clean, preprocess, and validate through a "
                    "reproducible Python pipeline.",
                    "Design a star-schema data warehouse in DuckDB.",
                    "Build engagement KPIs, course-outcome, and "
                    "engagement-pattern dashboards.",
                    "Train a supervised ML model for early dropout "
                    "prediction (first-30-day features).",
                    "Deploy free on Streamlit Community Cloud — public URL.",
                    "Validate accuracy through reconciliation and "
                    "document actionable insights for administrators.",
                ], size=18, spacing=6)
    add_page_number_footer(s, 5, total)
    add_speaker_notes(s,
        "Map objectives to project tasks (the 11 listed in your "
        "assignment). If they ask which task is not covered, say: "
        "all 11 are covered — objectives 1-3 map to collect/clean/"
        "warehouse; 4 maps to ETL; 5-7 map to dashboards; 8-11 map "
        "to validation and documentation.")

    # -------- slide 6: Scope
    s = blank_slide(prs)
    title_bar(s, 6, "Scope of the Project")
    # two-column layout: In scope / Out of scope
    add_textbox(s, Inches(0.7), Inches(1.4), Inches(5), Inches(0.5),
                "In scope", size=22, bold=True, colour=BRAND_DARK)
    add_bullets(s, Inches(0.7), Inches(2.0),
                Inches(6), Inches(5),
                [
                    "Descriptive analytics — KPIs, demographic "
                    "breakdowns.",
                    "Diagnostic analytics — peak/trough detection, "
                    "cross-category comparison.",
                    "Predictive analytics — early dropout-risk scoring.",
                    "Free public deployment on Streamlit Community Cloud.",
                ], size=17)
    add_textbox(s, Inches(7.2), Inches(1.4),
                Inches(5), Inches(0.5),
                "Out of scope", size=22, bold=True, colour=BRAND_DARK)
    add_bullets(s, Inches(7.2), Inches(2.0),
                Inches(6), Inches(5),
                [
                    "Real-time ingestion from a live LMS "
                    "(OULAD is a batch release).",
                    "Any personally identifiable information "
                    "(OULAD is fully anonymised).",
                    "Content-recommendation engines.",
                    "Multi-tenant or commercial deployment.",
                ], size=17)
    add_page_number_footer(s, 6, total)
    add_speaker_notes(s,
        "Clarity on scope prevents surprise questions. If asked "
        "'why not real-time?' — OULAD is a static release, real-time "
        "ingestion is listed in Chapter 8 as future work.")

    # -------- slide 7: Existing Systems
    s = blank_slide(prs)
    title_bar(s, 7, "Existing Systems — Where They Fall Short")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Native LMS dashboards (Moodle, Canvas, Blackboard) "
                    "— per-course only; no cross-course aggregation, "
                    "no predictive layer.",
                    "Commercial BI (Tableau, Power BI, Looker) — "
                    "licence costs ~$70/user/month; assume you already "
                    "have a warehouse.",
                    "Open-source BI (Superset, Metabase) — free but "
                    "general-purpose; no learning-analytics model, "
                    "no predictive layer.",
                    "None combine a domain data model + tailored "
                    "dashboards + native ML layer in one free deployable "
                    "artefact.",
                    "That gap is exactly what EduPulse fills.",
                ], size=19, spacing=8)
    add_page_number_footer(s, 7, total)
    add_speaker_notes(s,
        "Panel may ask why you didn't just use Tableau. Answer: "
        "licence cost and lack of a domain data model — you would "
        "still need to build the star schema and train the ML model "
        "yourself.")

    # -------- slide 8: Proposed System
    s = blank_slide(prs)
    title_bar(s, 8, "Proposed System — EduPulse",
              "Six tightly-integrated components on one free stack")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Reproducible ETL pipeline — download → clean → "
                    "load warehouse → train model.",
                    "DuckDB star-schema warehouse — 3 facts, 4 "
                    "dimensions, 3 analytical views.",
                    "scikit-learn Gradient Boosting classifier — "
                    "80.3% accuracy, 0.810 ROC AUC.",
                    "Streamlit multi-page dashboard — 5 purpose-built pages.",
                    "Free public deployment on Streamlit Community "
                    "Cloud — CD via GitHub push.",
                    "MIT-licensed — reproducible by any future "
                    "contributor.",
                ], size=19, spacing=8)
    add_page_number_footer(s, 8, total)
    add_speaker_notes(s,
        "Keep this slide crisp. The panel wants to know 'what did "
        "you build' — this is the summary. Emphasise 80.3% accuracy "
        "and 0.810 AUC — that number will come up again.")

    # -------- slide 9: Technology Stack
    s = blank_slide(prs)
    title_bar(s, 9, "Technology Stack")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Python 3.11 — one language across ETL, ML, and "
                    "dashboard.",
                    "DuckDB — embedded columnar OLAP warehouse, no "
                    "server needed.",
                    "pandas — tabular cleaning and feature engineering.",
                    "scikit-learn — GradientBoostingClassifier, "
                    "preprocessing pipeline, model persistence.",
                    "Streamlit — declarative Python web framework for "
                    "the dashboard.",
                    "Plotly — interactive charts (lines, bars, heatmaps).",
                    "Git + GitHub + Streamlit Community Cloud — "
                    "continuous deployment.",
                ], size=18, spacing=6)
    add_page_number_footer(s, 9, total)
    add_speaker_notes(s,
        "If asked 'why DuckDB instead of Postgres?' — DuckDB is "
        "columnar and OLAP-optimised; Postgres is row-store and "
        "OLTP-optimised. For analytical dashboards, DuckDB is "
        "10-100x faster on the same hardware.")

    # -------- slide 10: OULAD dataset
    s = blank_slide(prs)
    title_bar(s, 10, "Data Source — OULAD",
              "Open University Learning Analytics Dataset")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Released by Open University KMI (2014) under CC-BY-4.0.",
                    "32,593 students, 22 course presentations, 7 modules.",
                    "10.6 million click events in the Virtual Learning "
                    "Environment.",
                    "173,739 assessment submissions.",
                    "Complete demographic attributes (age band, region, "
                    "IMD band, prior education, disability).",
                    "Final outcome per enrolment: Pass / Fail / "
                    "Withdrawn / Distinction.",
                    "Downloaded automatically from UCI Machine Learning "
                    "Repository mirror.",
                ], size=18, spacing=6)
    add_page_number_footer(s, 10, total)
    add_speaker_notes(s,
        "OULAD numbers to remember: 32,593 students · 22 presentations "
        "· 10.6M clicks · 4 outcome classes. If asked about licence: "
        "CC-BY-4.0 permits redistribution with attribution.")

    # -------- slide 11: System Architecture (image)
    s = blank_slide(prs)
    title_bar(s, 11, "System Architecture",
              "Four-layer architecture from source to dashboard")
    fit_image(s, FIG / "figure_3_1_architecture.png",
              top=Inches(1.4), max_w=Inches(12.0), max_h=Inches(5.4))
    add_page_number_footer(s, 11, total)
    add_speaker_notes(s,
        "Walk left-to-right: Data source (OULAD) → ETL pipeline "
        "(4 Python modules) → Storage layer (DuckDB + joblib "
        "artefact) → Presentation layer (Streamlit). Each layer "
        "communicates with the next through a file artefact — no "
        "network protocols between layers.")

    # -------- slide 12: DFD Level 1 (image)
    s = blank_slide(prs)
    title_bar(s, 12, "Data Flow Diagram — Level 1",
              "Six processes, four data stores")
    fit_image(s, FIG / "figure_3_3_dfd_level_1.png",
              top=Inches(1.4), max_w=Inches(12.0), max_h=Inches(5.4))
    add_page_number_footer(s, 12, total)
    add_speaker_notes(s,
        "Processes 1.0-4.0 (Extract, Clean, Load, Train) are batch, "
        "run once at deployment. Processes 5.0 (Serve Dashboard) and "
        "6.0 (Score Cohort) run interactively at request time. Data "
        "stores D1-D4 are file artefacts.")

    # -------- slide 13: Star schema (image)
    s = blank_slide(prs)
    title_bar(s, 13, "Warehouse — Star Schema",
              "3 fact tables + 4 conformed dimensions")
    fit_image(s, FIG / "figure_4_5_star_schema.png",
              top=Inches(1.4), max_w=Inches(12.0), max_h=Inches(5.4))
    add_page_number_footer(s, 13, total)
    add_speaker_notes(s,
        "Facts in green at the centre: fact_engagement (10.6M rows), "
        "fact_assessment, fact_enrollment. Dimensions at the corners. "
        "dim_student is conformed across all three facts — that's "
        "what allows cross-fact queries on a single student key.")

    # -------- slide 14: Class Diagram (image)
    s = blank_slide(prs)
    title_bar(s, 14, "Class Diagram",
              "Python modules represented as UML classes")
    fit_image(s, FIG / "figure_4_1_class_diagram.png",
              top=Inches(1.4), max_w=Inches(12.0), max_h=Inches(5.4))
    add_page_number_footer(s, 14, total)
    add_speaker_notes(s,
        "Six logical modules. Downloader + Cleaner = ingestion. "
        "WarehouseBuilder = storage. DropoutModelTrainer = ML. "
        "DashboardApp + WarehouseAccess = presentation. Each module "
        "is one file in src/.")

    # -------- slide 15: ML Model overview
    s = blank_slide(prs)
    title_bar(s, 15, "Machine-Learning Model",
              "Early dropout prediction from first-30-day engagement")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Algorithm — scikit-learn GradientBoostingClassifier "
                    "(200 trees, max depth 3, learning rate 0.08).",
                    "Feature window — first 30 days of each course "
                    "presentation (enforced in v_student_features view).",
                    "Numeric features — total clicks, active days, "
                    "avg assessment score, studied credits, "
                    "prior attempts, submissions.",
                    "Categorical features — gender, region, IMD band, "
                    "age band, prior education, disability.",
                    "Training — stratified 80/20 split, "
                    "random_state = 42, deterministic re-runs.",
                    "Bundle — full sklearn Pipeline persisted via joblib "
                    "with metrics and feature importances.",
                ], size=18, spacing=6)
    add_page_number_footer(s, 15, total)
    add_speaker_notes(s,
        "If asked 'why Gradient Boosting?' — tree-based ensembles "
        "dominate this problem class (Hellas et al. 2018 systematic "
        "review). Interpretable via feature importances; not a black "
        "box.")

    # -------- slide 16: Model Performance (Confusion matrix image)
    s = blank_slide(prs)
    title_bar(s, 16, "Model Performance",
              "Held-out test set — 6,519 students")
    fit_image(s, FIG / "figure_6_1_confusion_matrix.png",
              top=Inches(1.5), max_w=Inches(6.7), max_h=Inches(5.2),
              centre=False, left=Inches(0.4))
    add_bullets(s, Inches(7.5), Inches(1.6),
                Inches(5.5), Inches(5),
                [
                    "Accuracy — 0.797",
                    "ROC AUC — 0.810",
                    "True negatives (Retained correctly) — 4,170",
                    "True positives (Withdrawn caught) — 1,023",
                    "Baseline dropout rate — 31.2%",
                    "AUC matches published OULAD results (Aljohani "
                    "et al. 2019: 0.80–0.85).",
                ], size=17, spacing=8)
    add_page_number_footer(s, 16, total)
    add_speaker_notes(s,
        "Panel loves numbers. Say: '0.810 ROC AUC, 79.7% accuracy on "
        "6,519 held-out students; baseline dropout rate is 31.2%, "
        "so accuracy alone is misleading — AUC is the honest metric.'")

    # -------- slide 17: Dashboard — Home + Overview
    s = blank_slide(prs)
    title_bar(s, 17, "Dashboard — Home & Overview")
    fit_image(s, SHOT / "screenshot_5_1_home.png",
              top=Inches(1.4), max_w=Inches(6.3), max_h=Inches(4.7),
              centre=False, left=Inches(0.4))
    fit_image(s, SHOT / "screenshot_5_2_overview.png",
              top=Inches(1.4), max_w=Inches(6.3), max_h=Inches(4.7),
              centre=False, left=Inches(6.9))
    add_textbox(s, Inches(0.4), Inches(6.25),
                Inches(6.3), Inches(0.5),
                "Home — brand mark + platform KPIs (28.8K students, "
                "39.6M clicks, 173.7K submissions)",
                size=11, colour=MUTED, italic=True)
    add_textbox(s, Inches(6.9), Inches(6.25),
                Inches(6.3), Inches(0.5),
                "Overview — outcome distribution, daily VLE click "
                "timeline, gender split",
                size=11, colour=MUTED, italic=True)
    add_page_number_footer(s, 17, total)
    add_speaker_notes(s,
        "Two live screenshots. If asked what a KPI card shows — "
        "students, presentations, enrollments, clicks, submissions. "
        "Overview page renders in under 2 seconds against the "
        "10.6M-row engagement fact.")

    # -------- slide 18: Dashboard — Course Analytics + Engagement
    s = blank_slide(prs)
    title_bar(s, 18, "Dashboard — Course Analytics & Engagement Patterns")
    fit_image(s, SHOT / "screenshot_5_3_course_analytics.png",
              top=Inches(1.4), max_w=Inches(6.3), max_h=Inches(4.7),
              centre=False, left=Inches(0.4))
    fit_image(s, SHOT / "screenshot_5_4_engagement_patterns.png",
              top=Inches(1.4), max_w=Inches(6.3), max_h=Inches(4.7),
              centre=False, left=Inches(6.9))
    add_textbox(s, Inches(0.4), Inches(6.25),
                Inches(6.3), Inches(0.5),
                "Course Analytics — enrolment popularity, pass/withdraw "
                "rates, module × presentation heatmap",
                size=11, colour=MUTED, italic=True)
    add_textbox(s, Inches(6.9), Inches(6.25),
                Inches(6.3), Inches(0.5),
                "Engagement Patterns — click and active-student "
                "timelines, peak/trough day tables",
                size=11, colour=MUTED, italic=True)
    add_page_number_footer(s, 18, total)
    add_speaker_notes(s,
        "Course Analytics answers 'which course is popular?' and "
        "'which is failing?' Engagement Patterns answers 'when do "
        "students engage most/least?' — filters by module and time "
        "window.")

    # -------- slide 19: Dashboard — Dropout Prediction
    s = blank_slide(prs)
    title_bar(s, 19, "Dashboard — Dropout Prediction",
              "The predictive layer administrators actually act on")
    fit_image(s, SHOT / "screenshot_5_5_dropout_prediction.png",
              top=Inches(1.4), max_w=Inches(10.0), max_h=Inches(5.4))
    add_page_number_footer(s, 19, total)
    add_speaker_notes(s,
        "This is the money slide. Model metrics at the top (0.803 "
        "accuracy, 0.810 AUC), feature-importance bar chart below. "
        "The admin selects a cohort, sees the top-20 at-risk "
        "students, and exports the full ranked list as CSV for "
        "outreach.")

    # -------- slide 20: Dashboard — Insights
    s = blank_slide(prs)
    title_bar(s, 20, "Dashboard — Insights & Recommendations",
              "Four validated findings, each with an admin action")
    fit_image(s, SHOT / "screenshot_5_6_insights.png",
              top=Inches(1.4), max_w=Inches(10.0), max_h=Inches(5.4))
    add_page_number_footer(s, 20, total)
    add_speaker_notes(s,
        "Four findings: (1) early engagement predicts retention, "
        "(2) withdrawal rate varies ~20 points across modules, "
        "(3) IMD deprivation band correlates with withdrawal, "
        "(4) prior attempts are a risk signal. Each with a "
        "concrete admin recommendation.")

    # -------- slide 21: Testing + Deployment
    s = blank_slide(prs)
    title_bar(s, 21, "Testing & Deployment")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Unit tests — 8 tests across cleaner, warehouse "
                    "loader, model trainer.",
                    "Integration test — full ETL pipeline; verifies "
                    "row counts, model AUC ≥ 0.75.",
                    "System tests — every dashboard page, filter, "
                    "and CSV export.",
                    "18 documented test cases, all PASS.",
                    "4 real bugs surfaced and fixed (OULAD URL, '?' "
                    "marker, INSERT column order, pyarrow wheel).",
                    "Deployment — public GitHub repo + Streamlit Cloud "
                    "auto-rebuild on every push to main.",
                    "Live URL — cost is zero, always-on, HTTPS via "
                    "Streamlit-managed TLS.",
                ], size=17, spacing=6)
    add_page_number_footer(s, 21, total)
    add_speaker_notes(s,
        "Panel will ask about bugs — call out Bug-004 (pyarrow → "
        "DuckDB parquet) because it's a real deployment war story. "
        "Bug-003 (SELECT column-order mismatch) was surfaced by "
        "your integration test catching an INSERT failure.")

    # -------- slide 22: Comparison + Results
    s = blank_slide(prs)
    title_bar(s, 22, "Results — Comparison with Existing Systems")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Cost — $0 per user per month (Tableau ≈ $70).",
                    "Free public deployment — Streamlit Community "
                    "Cloud (Tableau: no; Superset: possible w/ effort).",
                    "Learning-analytics data model — built-in "
                    "(Tableau / Superset: bring-your-own).",
                    "Native dropout-prediction layer — included "
                    "(others: custom build required).",
                    "Single-language stack — Python only (others: "
                    "Python + JS + SQL).",
                    "Operational surface area — one Streamlit process "
                    "(others: multi-tier).",
                    "Query latency — sub-200 ms on 10.6M-row fact table.",
                ], size=17, spacing=6)
    add_page_number_footer(s, 22, total)
    add_speaker_notes(s,
        "Only claim EduPulse wins on things it actually wins on. "
        "Acknowledge the trade-off: Tableau/Superset offer "
        "drag-and-drop authoring, EduPulse does not — but for a "
        "designed dashboard aimed at one persona, that trade is "
        "acceptable.")

    # -------- slide 23: Future Scope
    s = blank_slide(prs)
    title_bar(s, 23, "Future Scope",
              "Six concrete next steps — not vague generalities")
    add_bullets(s, Inches(0.7), Inches(1.5),
                Inches(12), Inches(5.6),
                [
                    "Real-time ingestion from live LMS "
                    "(Moodle / Canvas REST APIs).",
                    "Role-based access control + multi-tenant deployment.",
                    "Calibrated probabilities (isotonic regression) "
                    "+ per-student SHAP explanations.",
                    "Automated weekly outreach digests via GitHub "
                    "Actions cron.",
                    "Mobile-responsive layout for on-the-go "
                    "administrators.",
                    "Periodic model retraining + drift monitoring badge.",
                ], size=19, spacing=8)
    add_page_number_footer(s, 23, total)
    add_speaker_notes(s,
        "Every future-scope item is scoped and has an implementation "
        "hint — this is what the MCA guideline calls out as the "
        "difference between good and vague future-scope sections.")

    # -------- slide 24: Q&A / Thank you
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BRAND_DARK
    bg.line.fill.background()
    add_textbox(s, Inches(0.7), Inches(2.7),
                Inches(12), Inches(1.4),
                "Thank you", size=72, bold=True, colour=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(4.1),
                Inches(12), Inches(0.7),
                "Questions?",
                size=32, colour=WHITE, align=PP_ALIGN.CENTER,
                italic=True)
    add_textbox(s, Inches(0.7), Inches(5.5),
                Inches(12), Inches(0.4),
                "GitHub  ·  github.com/deepakjoshi-dj/edupulse",
                size=17, colour=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.7), Inches(5.9),
                Inches(12), Inches(0.4),
                "Deepak Kumar Joshi   ·   O24MCA112387   ·   "
                "MCA (CU Online)",
                size=15, colour=WHITE, align=PP_ALIGN.CENTER)
    add_speaker_notes(s,
        "Close politely: 'Thank you for your attention — I'm happy "
        "to take questions.' Sit down. Panel typically asks: (a) "
        "why this dataset, (b) why this algorithm, (c) how does it "
        "scale, (d) what would you change if you did it again. "
        "You have answers for all four.")

    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
